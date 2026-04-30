# src/scraper/bbtc_scraper.py
from pandas.core.tools.datetimes import to_datetime
import os
import sys
from pathlib import Path
import cloudscraper
import re
import unicodedata
from bs4 import BeautifulSoup
from datetime import datetime, timezone
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation

ROOT_DIR = Path(__file__).resolve().parents[2]
sys.path.insert(0, str(ROOT_DIR))

from src.storage.sqlite_store import SermonRegistry
from src.ingestion.file_classifier import classify_file
import urllib.parse

_RESOURCE_EXTENSIONS = (".pdf", ".pptx", ".ppt", ".docx", ".doc")
_LANG_SUFFIX = {"English": "audio-sermons", "Mandarin": "mw-sermons"}
_HEADERS = {
    "User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36"
}


class BBTCScraper:
    def __init__(
        self,
        download_dir: str = "data/sermons",
        staging_dir: str = "data/staging",
        registry: SermonRegistry | None = None,
    ):
        self._download_dir = download_dir
        self._staging_dir = staging_dir
        self._registry = registry or SermonRegistry()
        self._scraper = cloudscraper.create_scraper()
        os.makedirs(download_dir, exist_ok=True)
        os.makedirs(staging_dir, exist_ok=True)

    def _archive_url(self, year: int, lang: str) -> str:
        suffix = _LANG_SUFFIX.get(lang, "audio-sermons")
        return f"https://www.bbtc.com.sg/{suffix}-{year}/"

    def _extract_file_links_from_page(self, page_url: str) -> list[str]:
        try:
            response = self._scraper.get(page_url, headers=_HEADERS, timeout=15)
            response.raise_for_status()
            soup = BeautifulSoup(response.text, "html.parser")
            links = []
            for tag in soup.find_all("a", href=True):
                href = tag["href"]
                if any(href.lower().endswith(ext) for ext in _RESOURCE_EXTENSIONS):
                    if not href.startswith("http"):
                        base = page_url.rstrip("/")
                        href = f"{base}/{href.lstrip('/')}"
                    links.append(href)
            return list(set(links))
        except Exception as e:
            print(f"⚠️  Failed to get links from {page_url}: {e}")
            return []

    def _download_file(self, url: str, dest_path: str):
        response = self._scraper.get(url, headers=_HEADERS, timeout=30, stream=True)
        response.raise_for_status()
        with open(dest_path, "wb") as f:
            for chunk in response.iter_content(chunk_size=8192):
                f.write(chunk)

    def _clean_text(self, text: str) -> str:
        text = text.replace("\x00", "")
        text = "".join(
            ch for ch in text
            if unicodedata.category(ch)[0] != "C" or ch in "\n\t "
        )
        return re.sub(r"\s+", " ", text).strip()

    def _extract_text_from_file(self, file_path: str) -> tuple[str, str]:
        ext = os.path.splitext(file_path)[1].lower().lstrip(".")
        text = ""
        quality = "failed"

        try:
            if ext == "pdf":
                doc = fitz.open(file_path)
                text = " ".join(page.get_text("text", sort=True) for page in doc).strip()
                quality = "text" if len(text) > 20 else "failed"
                doc.close()
            elif ext in ("pptx", "ppt"):
                prs = Presentation(file_path)
                parts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            parts.append(shape.text)
                    if slide.has_notes_slide:
                        parts.append(slide.notes_slide.notes_text_frame.text)
                text = " ".join(parts).strip()
                quality = "text" if text else "failed"
            elif ext in ("docx", "doc"):
                doc = Document(file_path)
                text = " ".join(p.text for p in doc.paragraphs).strip()
                quality = "text" if text else "failed"
        except Exception as e:
            print(f"❌ Error extracting from {file_path}: {e}")
            return "", "failed"

        return self._clean_text(text), quality

    def _process_link(self, url: str, year: int, lang: str):
        if not self._registry.is_new(url):
            return

        basename = os.path.basename(url.split('?')[0])
        basename = urllib.parse.unquote(basename)
        filename = f"{lang}_{year}_{basename}"

        # Skip handout files before downloading
        if classify_file(filename) == "handout":
            print(f"⏭️  Skipping handout: {filename}")
            return

        staging_path = os.path.join(self._staging_dir, filename)

        # Download if not exists in staging
        if not os.path.exists(staging_path):
            try:
                self._download_file(url, staging_path)
                print(f"📥 Downloaded [{lang}/{year}]: {filename}")
            except Exception as e:
                print(f"⚠️  Failed to download {url}: {e}")
                return

        # Extract text
        print(f"🔍 Extracting [{lang}/{year}]: {filename}")
        extracted_text, quality = self._extract_text_from_file(staging_path)
        
        if quality == "failed":
            print(f"⚠️  Extraction failed for {filename}")
            status = "failed"
        else:
            # Save extracted text to a file in download_dir
            text_filename = os.path.splitext(filename)[0] + ".txt"
            text_path = os.path.join(self._download_dir, text_filename)
            with open(text_path, "w", encoding="utf-8") as f:
                f.write(extracted_text)
            status = "extracted"
            print(f"📝 Saved text: {text_filename}")

        # Update registry
        sermon_id = re.sub(r"\W+", "_", filename)
        record = {
            "sermon_id": sermon_id,
            "filename": filename,
            "url": url,
            "language": lang,
            "year": year,
            "status": status,
            "date_scraped": datetime.now(timezone.utc).isoformat(),
            "file_type": os.path.splitext(filename)[1].lstrip(".")
        }
        self._registry.insert_sermon(record)
        self._registry.mark_processed(url)

    def scrape_year(self, year: int, lang: str = "English"):
        url = self._archive_url(year, lang)
        try:
            response = self._scraper.get(url, headers=_HEADERS, timeout=15)
            response.raise_for_status()
            soup = BeautifulSoup(response.text, "html.parser")
            
            content = soup.find("div", class_="content")
            if not content:
                print(f"⚠️  Could not find content section on {url}")
                return
                
            sermon_links = set()
            direct_files = set()
            for a in content.find_all("a", href=True):
                href = a["href"]
                if not href.startswith("http"):
                    base = url.rstrip("/")
                    href = f"{base}/{href.lstrip('/')}"
                
                if "addtoany" in href or href == url or href == url + "/":
                    continue
                    
                parsed_path = urllib.parse.urlparse(href).path.lower()
                if any(parsed_path.endswith(ext) for ext in _RESOURCE_EXTENSIONS):
                    direct_files.add(href)
                elif href.startswith("https://www.bbtc.com.sg/"):
                    sermon_links.add(href)
            
            print(f"🔎 Found {len(sermon_links)} sermon pages and {len(direct_files)} direct files for {lang} {year}")
            
            # Process direct files first
            for file_url in direct_files:
                self._process_link(file_url, year=year, lang=lang)

            # Process sermon pages
            for sermon_url in sermon_links:
                file_links = self._extract_file_links_from_page(sermon_url)
                for link in file_links:
                    self._process_link(link, year=year, lang=lang)
                    
        except Exception as e:
            print(f"❌ Could not scrape {url}: {e}")

if __name__ == "__main__":
    import sys
    year = int(sys.argv[1]) if len(sys.argv) > 1 else datetime.now().year  
    print(f"🚀 Starting scraper for year {year}...")
    scraper = BBTCScraper()
    scraper.scrape_year(year)
    print("✅ Done!")
